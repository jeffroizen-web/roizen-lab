#!/usr/bin/env python3
"""
Extract images from PowerPoint presentations for the Roizen Lab website.
Saves images as PNGs with a manifest.json for each presentation.
"""

import json
import os
import sys
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu


def get_slide_text(slide):
    """Extract title and first text box content from a slide."""
    title = ""
    first_text = ""

    # First, grab the title directly if available
    try:
        if slide.shapes.title:
            title = slide.shapes.title.text.strip()
    except Exception:
        pass

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # Check if this is the title shape
            try:
                if slide.shapes.title and shape.shape_id == slide.shapes.title.shape_id:
                    title = text
                    continue
            except Exception:
                pass
            # Check placeholder
            try:
                pf = shape.placeholder_format
                if pf is not None and pf.idx == 0:
                    title = text
                    continue
            except (ValueError, AttributeError):
                pass
            # Otherwise it's body text
            if not first_text:
                first_text = text

    return title, first_text


def extract_images_from_pptx(pptx_path, output_dir, label=""):
    """Extract all images from a PowerPoint file."""
    print(f"\n{'='*60}")
    print(f"Processing: {label or pptx_path}")
    print(f"Output to:  {output_dir}")
    print(f"{'='*60}")

    if not os.path.exists(pptx_path):
        print(f"ERROR: File not found: {pptx_path}")
        return 0

    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    print(f"Total slides: {total_slides}")

    manifest = []
    image_count = 0
    skipped_small = 0

    def iter_shapes(shapes):
        """Recursively yield shapes, descending into group shapes."""
        for shape in shapes:
            yield shape
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                try:
                    yield from iter_shapes(shape.shapes)
                except Exception:
                    pass

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title, slide_text = get_slide_text(slide)
        img_on_slide = 0

        for shape in iter_shapes(slide.shapes):
            # Check if shape has an image
            if not shape.shape_type or not hasattr(shape, "image"):
                continue

            try:
                image = shape.image
                image_blob = image.blob
                content_type = image.content_type
            except Exception:
                continue

            # Determine file extension from content type
            ext_map = {
                'image/png': '.png',
                'image/jpeg': '.jpg',
                'image/gif': '.gif',
                'image/tiff': '.tiff',
                'image/bmp': '.bmp',
                'image/x-wmf': '.wmf',
                'image/x-emf': '.emf',
                'image/wmf': '.wmf',
                'image/emf': '.emf',
                'application/x-wmf': '.wmf',
                'application/x-emf': '.emf',
            }
            orig_ext = ext_map.get(content_type, '')

            # For WMF/EMF files, Pillow typically cannot load them properly.
            # Save as raw blob with original extension, and also try to get
            # dimensions from the shape itself.
            is_vector = orig_ext in ('.wmf', '.emf')

            if is_vector:
                # Use shape dimensions (in EMU) to estimate pixel size
                # 1 inch = 914400 EMU, assume 96 DPI for screen
                try:
                    width = int(shape.width / 914400 * 96) if shape.width else 0
                    height = int(shape.height / 914400 * 96) if shape.height else 0
                except Exception:
                    width, height = 0, 0

                if width < 50 or height < 50:
                    skipped_small += 1
                    continue

                img_on_slide += 1
                image_count += 1

                # Save as original vector format
                filename = f"slide{slide_idx:02d}_img{img_on_slide}{orig_ext}"
                filepath = os.path.join(output_dir, filename)
                with open(filepath, 'wb') as f:
                    f.write(image_blob)

            else:
                # Raster image — load with Pillow
                try:
                    img = Image.open(BytesIO(image_blob))
                except Exception as e:
                    print(f"  Slide {slide_idx}: Could not open image - {e}")
                    continue

                width, height = img.size

                # Skip very small images (icons, bullets, etc.)
                if width < 50 or height < 50:
                    skipped_small += 1
                    continue

                img_on_slide += 1
                image_count += 1

                # Build filename: slide03_img1.png
                filename = f"slide{slide_idx:02d}_img{img_on_slide}.png"
                filepath = os.path.join(output_dir, filename)

                # Convert to PNG and save
                try:
                    if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
                        img = img.convert('RGBA')
                    else:
                        img = img.convert('RGB')
                    img.save(filepath, 'PNG')
                except Exception as e:
                    # Fallback: save the raw blob with guessed extension
                    print(f"  Slide {slide_idx}: PNG conversion failed ({e}), saving raw blob")
                    fallback_ext = orig_ext if orig_ext else '.bin'
                    filename = f"slide{slide_idx:02d}_img{img_on_slide}{fallback_ext}"
                    filepath = os.path.join(output_dir, filename)
                    with open(filepath, 'wb') as f:
                        f.write(image_blob)

            # Add to manifest
            manifest.append({
                "filename": filename,
                "slide_number": slide_idx,
                "image_index_on_slide": img_on_slide,
                "width": width,
                "height": height,
                "original_content_type": content_type,
                "slide_title": slide_title,
                "slide_text": slide_text[:200] if slide_text else "",
            })

            print(f"  Slide {slide_idx:2d} -> {filename} ({width}x{height}) "
                  f"[{slide_title[:50] if slide_title else 'no title'}]")

    # Also extract images from the slide relationships (embedded media)
    # that might not be directly on shapes (e.g., background images)
    # We already capture shape-based images above; this is a safety net.

    # Save manifest
    manifest_path = os.path.join(output_dir, "manifest.json")
    manifest_data = {
        "source_file": os.path.basename(pptx_path),
        "source_path": pptx_path,
        "total_slides": total_slides,
        "total_images_extracted": image_count,
        "small_images_skipped": skipped_small,
        "images": manifest
    }

    with open(manifest_path, 'w') as f:
        json.dump(manifest_data, f, indent=2)

    print(f"\nSummary for {label or os.path.basename(pptx_path)}:")
    print(f"  Total slides:          {total_slides}")
    print(f"  Images extracted:      {image_count}")
    print(f"  Small images skipped:  {skipped_small}")
    print(f"  Manifest saved:        {manifest_path}")

    return image_count


def main():
    base_output = "/Users/roizenj/Desktop/Claude Apps/Roizen Lab/extracted-figures"

    presentations = [
        {
            "label": "Priority 1: 2026 Research Talk (87 slides)",
            "path": "/Users/roizenj/Documents/presentations and posters/02 2026 JH all inverted v 2.pptx",
            "output": os.path.join(base_output, "2026-talk"),
        },
        {
            "label": "Priority 2: Cool Science Pictures",
            "path": "/Users/roizenj/Documents/Personal/RoizenWebsite/science pics/cool science pictures.pptx",
            "output": os.path.join(base_output, "science-pics"),
        },
        {
            "label": "Priority 3: Lab Paradigm Figures",
            "path": "/Users/roizenj/Documents/Personal/RoizenWebsite/science pics/coplies_of_all_figures.pptx",
            "output": os.path.join(base_output, "paradigm-figures"),
        },
    ]

    grand_total = 0

    for pres in presentations:
        os.makedirs(pres["output"], exist_ok=True)
        count = extract_images_from_pptx(pres["path"], pres["output"], pres["label"])
        grand_total += count

    print(f"\n{'='*60}")
    print(f"GRAND TOTAL: {grand_total} images extracted across all presentations")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
